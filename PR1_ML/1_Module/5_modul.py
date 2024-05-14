import pandas as pd
import numpy as np
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.cluster import KMeans
from sklearn.decomposition import LatentDirichletAllocation
import umap
import matplotlib.pyplot as plt
import os
from pptx import Presentation
from pptx.util import Pt, Inches
from transformers import GPT2LMHeadModel, GPT2Tokenizer
import datetime

# 데이터 불러오기
path = "C:\\Users\\user\\OneDrive\\바탕 화면\\NEW\\3_ML_PM\\1_DB\\RAW_Shijin_tech.xlsx"
data = pd.read_excel(path)

# 결측치를 빈 문자열로 대체
data['요약(원문)'] = data['요약(원문)'].fillna('')

# 텍스트 데이터 전처리를 위한 Tfidf 벡터화
tfidf_vectorizer = TfidfVectorizer(max_features=1000, stop_words='english')
tfidf_matrix = tfidf_vectorizer.fit_transform(data['요약(원문)'])

# KMeans 클러스터링
num_clusters = 5
km = KMeans(n_clusters=num_clusters)
km.fit(tfidf_matrix)
clusters = km.labels_

# 클러스터 결과 추가
data['Cluster'] = clusters

# sklearn LDA 모델링
lda = LatentDirichletAllocation(n_components=5, max_iter=10, learning_method='online', random_state=0)
lda.fit(tfidf_matrix)

# UMAP을 이용한 시각화
reducer = umap.UMAP()
embedding = reducer.fit_transform(tfidf_matrix.toarray())

# 시각화 그래프 생성
plt.figure(figsize=(10, 6))
plt.scatter(embedding[:, 0], embedding[:, 1], c=data['Cluster'], cmap='Spectral', s=50)
plt.colorbar(label='Cluster')
plt.title('UMAP projection of patent dataset')
plt.xlabel('UMAP Dimension 1')
plt.ylabel('UMAP Dimension 2')
plt.text(0.5, 0.5, 'UMAP (Uniform Manifold Approximation and Projection)', horizontalalignment='center',
         verticalalignment='center', transform=plt.gca().transAxes, fontsize=12, bbox=dict(facecolor='white', alpha=0.5))

# UMAP 시각화 이미지 저장
umap_img_path = "C:\\Users\\user\\OneDrive\\바탕 화면\\NEW\\3_ML_PM\\3_Plot\\UMAP_projection_v1.png"
plt.savefig(umap_img_path)
plt.close()

# ChatGPT가 읽어야 할 데이터 텍스트 파일로 저장
sub_data_path = "C:\\Users\\user\\OneDrive\\바탕 화면\\NEW\\3_ML_PM\\6_sub_data\\data_for_gpt.txt"
data.to_csv(sub_data_path, header=None, index=None, sep='\t', mode='a')

# 결과 저장 
processed_data_path = "C:\\Users\\user\\OneDrive\\바탕 화면\\NEW\\3_ML_PM\\5_Processed\\Processed_Shijin_tech_v1.xlsx"
data.to_excel(processed_data_path)

# MS-PowerPoint 프레젠테이션 생성
prs = Presentation()

# 제목 슬라이드 추가
slide_title = prs.slides.add_slide(prs.slide_layouts[0])
title = slide_title.shapes.title
subtitle = slide_title.placeholders[1]
title.text = "특허 데이터 분석 보고서"
subtitle.text = "by KoGPT-2"

# 보고서 내용 생성을 위한 KoGPT-2 텍스트 생성 함수 정의
def generate_text(prompt):
    input_ids = tokenizer.encode(prompt, return_tensors="pt")
    output = model.generate(input_ids, max_length=150, num_return_sequences=1, early_stopping=True)
    return tokenizer.decode(output[0], skip_special_tokens=True)

# 모델과 토크나이저 로드
tokenizer = GPT2Tokenizer.from_pretrained("skt/kogpt2-base-v2", use_fast=False)
model = GPT2LMHeadModel.from_pretrained("skt/kogpt2-base-v2")

# LDA 모델링 결과 해석을 위한 KoGPT-2 텍스트 생성
lda_prompt = "특허 데이터에 대한 Latent Dirichlet Allocation (LDA) 모델링 결과를 해석하세요."
lda_summary = generate_text(lda_prompt)

# 내용 슬라이드 추가
slide_content = prs.slides.add_slide(prs.slide_layouts[1])
title, content = slide_content.shapes.title, slide_content.placeholders[1]
title.text = "보고서 내용"
content.text = lda_summary

# 시각화 슬라이드 추가
slide_visualization = prs.slides.add_slide(prs.slide_layouts[1])
title, content = slide_visualization.shapes.title, slide_visualization.placeholders[1]
title.text = "시각화 결과 해석"

# UMAP 이미지 해석을 위한 KoGPT-2 텍스트 생성
umap_prompt = "특허 데이터의 UMAP 시각화 결과를 해석하세요."
umap_summary = generate_text(umap_prompt)

# 이미지 파일이 존재하는지 확인하고 삽입
if os.path.exists(umap_img_path):
    # 이미지 추가
    pic = slide_visualization.shapes.add_picture(umap_img_path, Inches(0.5), Inches(1.5), width=Inches(9), height=Inches(5))
    content.text = umap_summary
else:
    content.text = "시각화 이미지를 찾을 수 없습니다."

# 폰트 크기 및 글꼴 설정
for slide in prs.slides:
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(20)
                    run.font.name = "맑은 고딕"

# 프레젠테이션 저장
ppt_save_path = "C:\\Users\\user\\OneDrive\\바탕 화면\\NEW\\3_ML_PM\\4_PPT"
ppt_files = os.listdir(ppt_save_path)
ppt_num = len(ppt_files) + 1
ppt_file_path = os.path.join(ppt_save_path, f"발표_v1_{ppt_num}.pptx")
prs.save(ppt_file_path)

print(f"발표 파일 저장 완료: {ppt_file_path}")
